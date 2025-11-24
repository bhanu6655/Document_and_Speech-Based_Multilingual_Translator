# document_loader.py
"""
Document and image loader for the Multilingual Translator.

Supports:
- .txt
- .docx
- .pdf
- .pptx
- image files (.png, .jpg, .jpeg, .bmp, .tiff) via OCR
"""

import os
from typing import List

from docx import Document
import pdfplumber
from pptx import Presentation
from PIL import Image
import pytesseract


TEXT_EXTENSIONS = {".txt"}
DOCX_EXTENSIONS = {".docx"}
PDF_EXTENSIONS = {".pdf"}
PPTX_EXTENSIONS = {".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".tif", ".tiff"}


def load_text_from_file(file_path: str) -> str:
    """
    Detect the file type by extension and extract text.
    Returns the extracted text as a single string.
    Raises ValueError for unsupported types.
    """
    ext = os.path.splitext(file_path)[1].lower()

    if ext in TEXT_EXTENSIONS:
        return _load_txt(file_path)
    elif ext in DOCX_EXTENSIONS:
        return _load_docx(file_path)
    elif ext in PDF_EXTENSIONS:
        return _load_pdf(file_path)
    elif ext in PPTX_EXTENSIONS:
        return _load_pptx(file_path)
    elif ext in IMAGE_EXTENSIONS:
        return _load_image_ocr(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


# === Individual loaders ===

def _load_txt(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _load_docx(file_path: str) -> str:
    doc = Document(file_path)
    content: List[str] = []

    # Paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            content.append(p.text)

    # Tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    content.append(cell.text)

    return "\n\n".join(content)  # separate paragraphs clearly


def _load_pdf(file_path: str) -> str:
    """
    Extract text from PDF page by page.
    """
    content: List[str] = []
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            if text.strip():
                content.append(text.strip())
    return "\n\n".join(content)


def _load_pptx(file_path: str) -> str:
    """
    Extract text from PPTX (slides).
    """
    prs = Presentation(file_path)
    content: List[str] = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    content.append(txt)

    return "\n\n".join(content)


def _load_image_ocr(file_path: str) -> str:
    """
    Use Tesseract OCR to extract text from an image.
    Make sure Tesseract is installed on the system.
    """
    img = Image.open(file_path)
    # You can pass lang parameter if you know image language, e.g. lang="eng"
    text = pytesseract.image_to_string(img)
    return text.strip()
